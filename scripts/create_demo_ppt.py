"""使用 python-pptx 直接生成演示文稿"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x0D, 0x11, 0x17)
CYAN = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x88, 0x99, 0xAA)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, font_name='Calibri', alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

# ====== Slide 1: Title ======
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BG)
add_textbox(slide1, 1.5, 1.8, 10, 1.5, "AI Agents", 60, CYAN, True, 'Georgia', PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 3.3, 10, 1, "The Future of Automation", 32, WHITE, False, 'Georgia', PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 4.5, 10, 0.8, "How intelligent agents are reshaping work, creativity, and decision-making", 16, MUTED, False, 'Calibri', PP_ALIGN.CENTER)

# ====== Slide 2: What Are AI Agents ======
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, LIGHT_BG)
add_textbox(slide2, 0.8, 0.5, 11, 1, "What Are AI Agents?", 40, DARK_TEXT, True, 'Georgia')
add_textbox(slide2, 0.8, 1.5, 5.5, 4.5,
    "AI agents are autonomous software entities that perceive their environment, "
    "make decisions, and take actions to achieve specific goals.\n\n"
    "Key Characteristics:\n"
    "- Autonomy: Operate without human intervention\n"
    "- Reactivity: Respond to environmental changes\n"
    "- Proactivity: Take initiative to achieve goals\n"
    "- Social Ability: Communicate with other agents",
    16, DARK_TEXT)

add_textbox(slide2, 7, 1.5, 5.5, 1.5, "Reactive", 24, CYAN, True, 'Georgia')
add_textbox(slide2, 7, 2.8, 5.5, 1.5, "Proactive", 24, CYAN, True, 'Georgia')
add_textbox(slide2, 7, 4.1, 5.5, 1.5, "Collaborative", 24, CYAN, True, 'Georgia')

# ====== Slide 3: Key Technologies ======
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, LIGHT_BG)
add_textbox(slide3, 0.8, 0.5, 11, 1, "Key Technologies", 40, DARK_TEXT, True, 'Georgia')

techs = [
    ("LLMs & Transformers", "Large Language Models provide\nreasoning and language understanding"),
    ("Tool Use & Function Calling", "Agents interact with external\nAPIs, databases, and services"),
    ("Memory & Context", "Vector stores and context windows\nenable persistent knowledge"),
    ("Multi-Agent Systems", "Specialized agents collaborate\nto solve complex problems"),
]
for i, (title, desc) in enumerate(techs):
    x = 0.8 + (i % 2) * 6
    y = 1.8 + (i // 2) * 2.6
    add_textbox(slide3, x, y, 2.5, 0.5, title, 18, CYAN, True, 'Georgia')
    add_textbox(slide3, x, y + 0.5, 5, 1.2, desc, 14, DARK_TEXT)

# ====== Slide 4: Use Cases ======
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, LIGHT_BG)
add_textbox(slide4, 0.8, 0.5, 11, 1, "Real-World Use Cases", 40, DARK_TEXT, True, 'Georgia')

cases = [
    ("Customer Service", "AI agents handle 80% of\nroutine support tickets"),
    ("Code Generation", "Agents write, test, and\ndeploy software autonomously"),
    ("Data Analysis", "Multi-agent systems uncover\ninsights from complex datasets"),
]
for i, (title, stat) in enumerate(cases):
    x = 0.8 + i * 4.2
    add_textbox(slide4, x, 2.0, 4, 1, title, 22, CYAN, True, 'Georgia')
    add_textbox(slide4, x, 3.2, 4, 1, stat, 16, DARK_TEXT)

# ====== Slide 5: Conclusion ======
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, DARK_BG)
add_textbox(slide5, 1.5, 1.5, 10, 1.5, "2025+", 72, CYAN, True, 'Georgia', PP_ALIGN.CENTER)
add_textbox(slide5, 1.5, 3.0, 10, 0.8, "The Year of Agentic AI", 20, MUTED, False, 'Calibri', PP_ALIGN.CENTER)
add_textbox(slide5, 1.5, 4.2, 10, 1.5,
    'AI agents will assist us, collaborate with us, and\n'
    'operate autonomously at a scale we have never seen.', 22, WHITE, False, 'Georgia', PP_ALIGN.CENTER)

output_path = os.path.join(os.path.dirname(__file__), "..", "demo_ai_agents.pptx")
prs.save(output_path)
print(f"PPTX saved to: {os.path.abspath(output_path)}")
